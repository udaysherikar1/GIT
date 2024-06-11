
from pptx import Presentation
from pptx.util import Inches,Pt
from pptx.enum.shapes import MSO_SHAPE
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

import csv

def read_csv(file_path):
    data = []
    with open(file_path, 'r') as file:
        csv_reader = csv.reader(file)
        for row in csv_reader:
            data.append(row)
    return data
    

def Add_text (ppt_template,textInfo,left,top,width1,height1,font_size):
       
    # Create a presentation object
    prs = Presentation(ppt_template)
    # Get the last slide..
    last_slide = prs.slides[-1]
    
    # Add text box to the slide
    left = Inches(left)
    top = Inches(top)
    text_box = last_slide.shapes.add_textbox(left, top, width=Inches(width1), height=Inches(height1))
    
    # Add text to the text box
    text_frame = text_box.text_frame
    p = text_frame.add_paragraph()
    p.text = textInfo
    #p.font.bold = "bold"
    #p.font.italic = "italic"
    p.font.size = Pt(font_size)
    p.font.name = "Verdana"
    p.font.color.rgb = RGBColor(255,0,0)
    #p.font.alignment = "align"

    # Save the presentation
    prs.save(ppt_template)
    


def AddSlide (ppt_template,slideTitle):

    #open existing ppt
    prs = Presentation(ppt_template)
    # Add a new slide (you can specify the layout using the 'layout' parameter)
    slide_layout = prs.slide_layouts[6]  # Choose the layout you want, 0 is usually the title slide
    slide = prs.slides.add_slide(slide_layout)

    #title
    title_shape = slide.shapes.title
    title_shape.text = slideTitle

    # Save the modified presentation
    prs.save(ppt_template)
    print("Slide added successfully to", ppt_template)

    return slide

def AddImagesInLastSlide(ppt_template,image_path, l_inch, t_inch, w_inch, h_inch):

    #black color
    border_color = (0, 0, 0)
    border_width = 0.5

    prs = Presentation(ppt_template)
    # Get the last slide..
    last_slide = prs.slides[-1]
    
    left_inch = Inches(l_inch)           # Adjust the left position of the image
    top_inch = Inches(t_inch)            # Adjust the top position of the image
    width_inch = Inches(w_inch)          # Adjust the width of the image (optional)
    height_inch = Inches(h_inch)         # Adjust the height of the image (optional)

    # Add the image to the slide
    shape=last_slide.shapes.add_picture(image_path, left_inch, top_inch, width=width_inch, height=height_inch)
    picture = shape._pic
    outline = shape.line
    outline.color.rgb = RGBColor(*border_color)
    outline.width = Pt(border_width)

    prs.save(ppt_template)


def AddTable(ppt_template,n_rows, n_columns,l_inch, t_inch, table_width, table_height):

    prs = Presentation(ppt_template)
    # Get the last slide..
    last_slide = prs.slides[-1]
    
    left_inch = Inches(l_inch)           
    top_inch = Inches(t_inch)            
    width_inch = Inches(table_width)
    height_inch = Inches(table_height) 

    # Add the table to the slide
    table = last_slide.shapes.add_table(n_rows, n_columns, left_inch, top_inch, table_width, table_height).table
    
    
    # Set the column widths
    col_width = table_width / n_columns
    for col in range(n_columns):
        #table.columns[col].width = Inches(col_width / Inches(1))
        table.columns[col].width = Inches(1)

    # Set the row heights
    row_height = table_height / n_rows
    for row in range(n_rows):
        #table.rows[row].height = Inches(row_height / Inches(1))
        table.rows[row].height = Inches(0.5)

    # # Populate the table with some data
    # for row in range(n_rows):
        # for col in range(n_columns):
            # cell = table.cell(row, col)
            # cell.text = f'Row {row+1}, Col {col+1}'

    # Save the presentation
    prs.save(ppt_template)